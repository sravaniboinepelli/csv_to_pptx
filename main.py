import argparse
from pptx import Presentation

from data import get_data_csv
from utils import duplicate_slide, add_images, replace_text_placeholders

ap = argparse.ArgumentParser(
        prog="csv2pptx",
        description="Generate slides in a presentation using a "
                    "template with placeholders"
        )

ap.add_argument("-t", "--template", required=True,
        help="Path to the template presentation")
ap.add_argument("-d", "--data", required=True,
        help="Path to the csv data")

args = vars(ap.parse_args())

pr = Presentation(pptx=args["template"])
header, rows = get_data_csv(args["data"])

for row in rows:
    index = 0
    if "{{index}}" in header:
        index = int(row[header.index("{{index}}")])

    new_slide = duplicate_slide(pr, index)

    # add all images
    success_img = add_images(new_slide, header, row)
    print(f"Added {success_img} images")

    # replace all text placeholders
    success_txt = replace_text_placeholders(new_slide, header, row)
    print(f"Replaced {success_txt} placeholders")


pr.save("gen.pptx")
print("Done")
