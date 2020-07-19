# This file contains stuff related to actions on slides

import re
from copy import deepcopy
from pptx.util import Inches, Pt, Cm, Mm

from config import IMG_CONF

def is_text_placeholder(shape):
    """
    Checks whether this shape is a text placeholder or not

    Parameters:
        shape: A shape(object) in a slide [eg: textbox, image, etc]
    """

    # text placeholders look like {{<field>}}
    re_placeholder = re.compile(r"{{[^}]*}}")

    if shape.has_text_frame:
        text = shape.text_frame.text
        return re_placeholder.match(text) is not None

    return False


def duplicate_slide(pr,index=0):
    """
    Duplicates a slide by copying element by element
    The new slide **is added** at the end of presentation

    Stolen from https://stackoverflow.com/a/51154583 :P
    Some parts were modified

    Parameters:
        pr: the presentation object
        index: the index of the slide which has to be copied
                (defaults to the first slide)
    """

    template = pr.slides[index]

    # 0 corresponds to BLANK_SLIDE
    # I had to find this by printing each layout's name
    # The docs weren't good enough ¯\_(ツ)_/¯
    blank_slide_layout = pr.slide_layouts[0]

    copied_slide = pr.slides.add_slide(blank_slide_layout)

    for shape in template.shapes:
        el = shape.element
        newel = deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, "p:extLst")

    for _, value in template.part.rels.items():
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(value.reltype,
                    value._target, value.rId)

    return copied_slide


def replace_text_keep_formatting(text_frame, search_str, replace_str):
    """
    Simply changing text_frame.text doesn't keep formatting
    To keep formatting intact, we need to change the text at the lowest level

    Stolen from https://stackoverflow.com/a/56740560 :P
    Some parts were modified

    Parameters:
        text_frame: a ppt textbox basically
        search_str, replace_str: replace search_str with replace_str
    """

    text = text_frame.paragraphs[0].runs[0].text
    new_text = text.replace(search_str, replace_str)

    text_frame.paragraphs[0].runs[0].text = new_text


def add_images(slide, header, data):
    """
    Adds all images specified in csv
    For this to happen, you also have to update config.py
    The config for images is specified there seperately

    Parameters:
        slide: The slide in which images have to be added
        header: The csv header which contains image info
        data: The csv row which has file_path for images

    Returns:
        success: number of images successfully added
    """

    # string to function map
    unit_map = {
        "Inches": Inches,
        "Pt": Pt,
        "Cm": Cm,
        "Mm": Mm
    }

    success = 0

    for conf in IMG_CONF:
        name = conf.get("name")
        unit = unit_map[conf.get("unit")]

        try:
            img_index = header.index(name)
        except ValueError:
            print(f"The image placeholder {name} doesn't "
                    "exist in the csv")
            continue

        file_path = data[img_index]

        if file_path == '':
            continue

        left = unit(conf.get("left"))
        top = unit(conf.get("top"))
        width = unit(conf.get("width"))
        height = unit(conf.get("height"))

        slide.shapes.add_picture(file_path, left, top, width, height)

        success += 1

    return success


def replace_text_placeholders(slide, header, data):
    """
    Replaces all text placeholders within a slide

    Parameters:
        slide: The slide in which placeholder text is to be replaced
        header: The csv header which contains placeholder info
        data: The data which the placeholder is to be replaced with
    Returns:
        success: number of placeholders successfully replaced
    """

    success = 0

    for shape in slide.shapes:
        if is_text_placeholder(shape):
            placeholder = shape.text_frame.text

            try:
                data_index = header.index(placeholder)
            except ValueError:
                print(f"The placeholder {placeholder} doesn't "
                        "exist in the template")
                continue

            replace_text_keep_formatting(shape.text_frame,
                    placeholder, data[data_index])

            success += 1

    return success

